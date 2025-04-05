from fastapi import FastAPI, File, UploadFile, HTTPException, Form, Depends
from fastapi.responses import JSONResponse
from fastapi.security import HTTPBasic, HTTPBasicCredentials
import shutil
import os
import pdfplumber
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from dotenv import load_dotenv
from langchain_huggingface import HuggingFaceEndpoint
import secrets

# Load environment variables
load_dotenv()
token = os.getenv("HUGGINGFACEHUB_API_TOKEN")

app = FastAPI()
UPLOAD_DIR = "uploaded_docs"
os.makedirs(UPLOAD_DIR, exist_ok=True)

documents = {}
security = HTTPBasic()

# Mock user database
users_db = {
    "admin": {"password": "admin123", "role": "admin"},
    "user1": {"password": "user123", "role": "user"},
}

# ---------------------- Authentication ----------------------
def authenticate(credentials: HTTPBasicCredentials = Depends(security)):
    username = credentials.username
    password = credentials.password

    user = users_db.get(username)
    if not user or not secrets.compare_digest(user["password"], password):
        raise HTTPException(status_code=401, detail="Authentication failed.")
    return {"username": username, "role": user["role"]}

# ---------------------- Hugging Face Model ----------------------
llm = HuggingFaceEndpoint(
    repo_id="mistralai/Mixtral-8x7B-Instruct-v0.1",
    task="text-generation",
    huggingfacehub_api_token=token,
    temperature=0.7,
    max_new_tokens=300,
)

# ---------------------- Text Extraction ----------------------
def extract_text(file_path, file_type):
    try:
        if file_type == "pdf":
            with pdfplumber.open(file_path) as pdf:
                return "\n".join([page.extract_text() for page in pdf.pages if page.extract_text()])
        elif file_type == "docx":
            doc = Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs])
        elif file_type == "txt":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as file:
                return file.read()
        elif file_type == "xlsx":
            workbook = load_workbook(file_path)
            text = ""
            for sheet in workbook.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    text += " ".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
            return text
        elif file_type == "pptx":
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        else:
            return "Unsupported file type."
    except Exception as e:
        return f"Error extracting text: {str(e)}"

# ---------------------- Routes ----------------------
@app.get("/")
def read_root():
    return {"message": "Welcome to the Document Search Bot API"}

@app.post("/get-role/")
def get_role(credentials: HTTPBasicCredentials = Depends(security)):
    username = credentials.username
    password = credentials.password
    user = users_db.get(username)
    if not user or not secrets.compare_digest(user["password"], password):
        raise HTTPException(status_code=401, detail="Invalid credentials")
    return {"username": username, "role": user["role"]}

@app.post("/upload/")
def upload_file(file: UploadFile = File(...), user: dict = Depends(authenticate)):
    if user["role"] != "admin":
        raise HTTPException(status_code=403, detail="Only admins can upload files.")

    if file.size is not None and file.size > 2 * 1024 * 1024:
        raise HTTPException(status_code=413, detail="File size exceeds 2MB limit.")
    
    file_path = os.path.join(UPLOAD_DIR, file.filename)
    with open(file_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    file_extension = file.filename.split(".")[-1].lower()
    extracted_text = extract_text(file_path, file_extension)
    documents[file.filename] = extracted_text

    return {"message": f"File '{file.filename}' uploaded and processed successfully!"}

@app.get("/list-files/")
def list_files(user: dict = Depends(authenticate)):
    return {"uploaded_files": list(documents.keys())}

@app.delete("/delete/{filename}")
def delete_file(filename: str, user: dict = Depends(authenticate)):
    if user["role"] != "admin":
        raise HTTPException(status_code=403, detail="Only admins can delete files.")
    if filename not in documents:
        raise HTTPException(status_code=404, detail="File not found")
    del documents[filename]
    return {"message": f"File '{filename}' deleted successfully!"}

@app.post("/query/")
def query_document(query: str = Form(...), user: dict = Depends(authenticate)):
    if not documents:
        raise HTTPException(status_code=400, detail="No documents available for search")

    # Limit document content to avoid token overflow
    max_chunks = 3  # Tune this based on expected document size
    selected_docs = list(documents.values())[:max_chunks]
    context = "\n\n".join(selected_docs)

    prompt = f"""Use the following context to answer the question.
If the answer is not found, respond with 'Answer not found in the documents.'

Context:
{context}

Question: {query}
Answer:"""

    try:
        response = llm.invoke(prompt)
        return JSONResponse(content={"query": query, "answer": response})
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error from Hugging Face model: {str(e)}")
